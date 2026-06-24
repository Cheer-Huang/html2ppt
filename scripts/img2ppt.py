#!/usr/bin/env python3
"""
图片转 PPT 脚本
将 PNG 图片序列组装为 PowerPoint 演示文稿。

用法:
    python img2ppt.py <input_dir_or_files> [output.pptx] [--width 1280] [--height 720]

参数:
    input_dir_or_files: 图片目录或逗号分隔的图片文件列表
    output.pptx: 输出 PPT 路径（默认为 slides.pptx）
    --width: 幻灯片宽度像素（默认 1280，对应 16:9）
    --height: 幻灯片高度像素（默认 720）
"""

import sys
import os
import glob
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Emu
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("错误: 请先安装 python-pptx: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


# 像素到 EMU 的转换（96 DPI）
# 1 inch = 96 px = 914400 EMU
PX_TO_EMU = 914400 / 96


def images_to_pptx(image_files, output_path, slide_width_px=1280, slide_height_px=720):
    """
    将图片列表组装为 PPTX 文件。
    
    参数:
        image_files: 图片文件路径列表（按顺序排列）
        output_path: 输出 PPTX 文件路径
        slide_width_px: 幻灯片宽度（像素）
        slide_height_px: 幻灯片高度（像素）
    """
    if not image_files:
        print("错误: 没有找到图片文件", file=sys.stderr)
        sys.exit(1)
    
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片尺寸（像素转 EMU）
    prs.slide_width = Emu(int(slide_width_px * PX_TO_EMU))
    prs.slide_height = Emu(int(slide_height_px * PX_TO_EMU))
    
    # 使用空白布局
    blank_layout = prs.slide_layouts[6]  # 6 = Blank
    
    print(f"幻灯片尺寸: {slide_width_px}x{slide_height_px}px (16:9)")
    print(f"共 {len(image_files)} 张图片，开始组装 PPT...\n")
    
    for i, img_path in enumerate(image_files):
        if not os.path.exists(img_path):
            print(f"  [{i+1}] 跳过（文件不存在）: {img_path}")
            continue
        
        slide = prs.slides.add_slide(blank_layout)
        
        # 图片填满整个幻灯片
        slide.shapes.add_picture(
            img_path,
            left=Emu(0),
            top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )
        
        print(f"  [{i+1}/{len(image_files)}] {os.path.basename(img_path)} ✓")
    
    # 保存
    prs.save(output_path)
    print(f"\n完成！已保存至: {output_path}")
    print(f"文件大小: {os.path.getsize(output_path) / 1024 / 1024:.1f} MB")


def main():
    parser = argparse.ArgumentParser(description="将图片序列转为 PPTX")
    parser.add_argument("input", help="图片目录或逗号分隔的文件列表")
    parser.add_argument("output", nargs="?", default="slides.pptx", help="输出 PPTX 路径")
    parser.add_argument("--width", type=int, default=1280, help="幻灯片宽度像素（默认1280）")
    parser.add_argument("--height", type=int, default=720, help="幻灯片高度像素（默认720）")
    
    args = parser.parse_args()
    
    # 收集图片文件
    if os.path.isdir(args.input):
        # 目录模式：按文件名排序
        extensions = ["*.png", "*.PNG", "*.jpg", "*.JPG", "*.jpeg", "*.JPEG"]
        image_files = []
        for ext in extensions:
            image_files.extend(glob.glob(os.path.join(args.input, ext)))
        image_files.sort()
    elif "," in args.input:
        # 逗号分隔的文件列表
        image_files = [f.strip() for f in args.input.split(",")]
    else:
        # 单个文件
        image_files = [args.input]
    
    if not image_files:
        print(f"错误: 在 {args.input} 中未找到图片文件", file=sys.stderr)
        sys.exit(1)
    
    print(f"输入: {args.input}")
    print(f"输出: {args.output}")
    print(f"找到 {len(image_files)} 张图片\n")
    
    images_to_pptx(image_files, args.output, args.width, args.height)


if __name__ == "__main__":
    main()
